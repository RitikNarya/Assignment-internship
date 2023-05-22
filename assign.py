from docx import Document
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt
from pptx.dml.color import RGBColor
from mathpix.mathpix import MathPix

# Function to extract equations using Mathpix OCR
def extract_equation_from_image(image_path, app_id, app_key):
    mp = MathPix(app_id=app_id, app_key=app_key)
    return mp.image_to_latex(image_path)

# Function to add equations to PowerPoint slides
def add_equation_to_slide(prs, equation, slide_layout):
    slide = prs.slides.add_slide(slide_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = equation
    p.space_before = Pt(24)
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 0, 0)

# Load the .docx file
doc = Document('input.docx')

# Create a PowerPoint presentation
prs = Presentation()

# Set slide layout for equation slides
slide_layout = prs.slide_layouts[5]  # Use the desired layout index (5 represents the blank layout)

# Iterate through the document and process equations and images
for element in doc.element.body:
    if element.tag.endswith(('equation', 'pict')):
        if element.tag.endswith('equation'):
            # Extract equation from XML
            equation = element[2][0].text
            add_equation_to_slide(prs, equation, slide_layout)
        elif element.tag.endswith('pict'):
            # Save the image as 'equation.png' and extract equation using Mathpix OCR
            element[0][0].save('equation.png')
            equation = extract_equation_from_image('equation.png', 'YOUR_MATHPIX_APP_ID', 'YOUR_MATHPIX_APP_KEY')
            add_equation_to_slide(prs, equation, slide_layout)

# Save the PowerPoint presentation
prs.save('output.pptx')

